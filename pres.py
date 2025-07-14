import os
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageEnhance, ImageFilter
import subprocess

# CONFIGURATION
PIXABAY_API_KEY = "51175310-75e88d74fff8a9e10acf9cf0e"
SEARCH_TERMS = [
    "data science dark", "customer segmentation", "shopping mall", 
    "machine learning", "analytics dashboard", "urban lifestyle", 
    "target marketing", "business strategy", "innovation technology"
]
IMAGE_FOLDER = r"C:\Users\paulg\Mall Segmentation\images"
PPTX_PATH = r"C:\Users\paulg\Mall Segmentation\Mall_Customer_Segmentation_with_Backgrounds.pptx"
SLIDE_WIDTH, SLIDE_HEIGHT = Inches(13.33), Inches(7.5)  # 1920x1080 aspect

SLIDE_CONTENT = [
    ("Mall Customer Segmentation App",
     "Using KMeans Clustering to Identify Target Marketing Segments\nPresented by: [Your Group Name]\nDate: July 2025\nInstitution: [Insert Logo]"),
    ("Background / Case Study",
     "Client: UrbanEdge Living – operates furniture stores in urban malls\nBusiness Challenge:\n- Distinguish customer types\n- Launch targeted promotions\n- Personalize marketing campaigns\nObjective: Use machine learning to cluster customers based on data"),
    ("Dataset Description",
     "Source: Kaggle Mall Customer Dataset\nSize: 200 rows, 5 columns\nKey Features: Age, Gender, Annual Income, Spending Score"),
    ("Data Cleaning & Preprocessing",
     "Dropped CustomerID (non-informative)\nEncoded Gender: Male=0, Female=1\nScaled features using StandardScaler\nNo missing values found"),
    ("Modeling Approach",
     "Algorithm: KMeans Clustering\nValidation: Elbow Method selected k=5 clusters\nVisual confirmation of clusters"),
    ("Application Demo",
     "Platform: Python & Streamlit\nFeatures:\n- Enter new customer info\n- Predict customer segment\n- Recommend marketing strategy\n- Visualize customer among clusters"),
    ("Segment Profiles",
     "Cluster 0: Low Income & Low Score → Low engagement group\nCluster 1: High Income & High Score → VIP shoppers\nCluster 2: Mid Income & High Score → Trendy young spenders\nCluster 3: Mid Income & Mid Score → Average customers\nCluster 4: Low Income & High Score → Bargain seekers"),
    ("Business Conclusion",
     "Achievements:\n- Customer segmentation for smarter targeting\n- Personalized campaigns\n- Avoided ineffective mass marketing\nFuture Ideas:\n- Real-time POS input\n- Feedback metrics\n- Loyalty program integration"),
    ("Q&A / Thank You",
     "Thank you for your attention\nQuestions and feedback welcome\nTeam: [Your Names]\nCourse / Instructor: [Insert Here]")
]

SLIDE_NOTES = [
    "Introduce the project and its relevance to marketing teams.",
    "Explain the client background and why segmentation matters.",
    "Walk through key dataset features: age, gender, income, spending score.",
    "Discuss preprocessing steps: encoding, scaling, and outlier removal.",
    "Present KMeans clustering approach and elbow method for k determination.",
    "Showcase the Streamlit app functionality and user interface.",
    "Highlight segment characteristics and marketing insights.",
    "Summarize business outcomes and future roadmap.",
    "Thank the audience and open for Q&A."
]

# Ensure image folder exists
os.makedirs(IMAGE_FOLDER, exist_ok=True)


def download_and_process_image(term, slide_num):
    img_path = os.path.join(IMAGE_FOLDER, f"slide{slide_num}.jpg")
    if os.path.exists(img_path):
        print(f"✅ Cached image found for Slide {slide_num}")
        return img_path

    print(f"⬇️ Downloading HD image for Slide {slide_num}...")
    url = f"https://pixabay.com/api/?key={PIXABAY_API_KEY}&q={term}&image_type=photo&orientation=horizontal&category=business&colors=black&per_page=5"
    response = requests.get(url)
    results = response.json().get('hits', [])

    if not results:
        print(f"⚠️ No results for '{term}', using fallback dark background.")
        return None

    # Pick the most vibrant image (highest contrast)
    best_image_url = None
    max_contrast = 0
    for hit in results:
        image_url = hit['largeImageURL']
        img_resp = requests.get(image_url)
        img = Image.open(BytesIO(img_resp.content)).convert('RGB')
        contrast = ImageEnhance.Contrast(img).enhance(2).getextrema()[1][1] - img.getextrema()[0][0]
        if contrast > max_contrast:
            max_contrast = contrast
            best_image_url = image_url

    if best_image_url:
        img_resp = requests.get(best_image_url)
        img = Image.open(BytesIO(img_resp.content)).convert('RGB')
        img = img.resize((1920, 1080), Image.LANCZOS)
        img = img.filter(ImageFilter.GaussianBlur(3))  # Slight blur for text legibility
        img.save(img_path)
        print(f"✅ Saved: {img_path}")
        return img_path
    else:
        print(f"⚠️ Failed to download image for Slide {slide_num}")
        return None


def set_background(slide, img_path):
    if img_path:
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=SLIDE_WIDTH, height=SLIDE_HEIGHT)


def add_text(slide, title, content):
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.font.name = "Segoe UI"
    p.font.shadow = True  # Add shadow for visibility

    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(6), Inches(12), Inches(1.5))
    tf2 = content_box.text_frame
    p2 = tf2.add_paragraph()
    p2.text = content
    p2.font.size = Pt(30)
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.font.name = "Calibri Light"
    p2.font.shadow = True


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for i, (title, content) in enumerate(SLIDE_CONTENT, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        img_path = download_and_process_image(SEARCH_TERMS[i-1], i)
        set_background(slide, img_path)
        add_text(slide, title, content)

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = SLIDE_NOTES[i-1]

    prs.save(PPTX_PATH)
    print(f"\n🎉 Presentation saved at: {PPTX_PATH}")
    open_pptx(PPTX_PATH)


def open_pptx(path):
    try:
        os.startfile(path)  # Windows
    except AttributeError:
        subprocess.call(['open', path])  # macOS/Linux


if __name__ == "__main__":
    create_presentation()
